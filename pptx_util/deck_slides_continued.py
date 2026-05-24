"""Slides 7-11 for the Telemetry & Team Velocity deck."""

from .deck_primitives import (
    ACC,
    BODY_FG,
    MUTED_FG,
    OK_BG,
    OK_FG,
    ERR_BG,
    ERR_FG,
    chrome,
    tbl,
    call,
    txt,
    card,
)
from .deck_slides import TS

from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


def s7(p):
    s = p.slides.add_slide(p.slide_layouts[6])
    chrome(
        s,
        7,
        TS,
        "Grafana: Team Velocity Dashboard",
        "25 panels · 9 sections · 3 dropdown filters (user, lane, group) · Auto-loaded by /scaffold",
    )
    panels = [
        ("Turns by\nUser", "Bar chart"),
        ("Success vs\nFailure", "Donut"),
        ("Agent/Tool\nWorkload", "Bar chart"),
        ("Lane\nUsage", "Pie chart"),
        ("Iteration\nper Group", "Stat"),
        ("Active\nStories", "Table"),
        ("Pending\nReviews", "Stat"),
        ("Effort per\nGroup", "Bar gauge"),
        ("Effort per\nStory", "Bar gauge"),
        ("Success per\nAgent/Tool", "Bar gauge"),
        ("Exec Mode\nDistrib.", "Pie chart"),
        ("Turns by\nUser+Lane", "Bar gauge"),
        ("Command\nInvocations", "Bar gauge"),
        ("Tool\nActivity", "Bar gauge"),
        ("Sessions", "Stat"),
        ("Cost\n(USD)", "Stat"),
        ("Tokens\nUsed", "Stat"),
        ("Active\nTime", "Stat"),
        ("Developer\nOutput", "Bar gauge"),
        ("Native\nTokens", "Bar gauge"),
        ("Native\nCost", "Bar gauge"),
        ("Edit\nDecisions", "Bar gauge"),
        ("Skill\nUsage", "Bar gauge"),
        ("Installed\nSkills", "Table"),
        ("All Agent\nRuns", "Table"),
    ]
    xs, ys = 0.45, 1.22
    cw, rh = 1.75, 0.90
    g = 0.10
    for i, (t, d) in enumerate(panels):
        x = xs + (i % 5) * (cw + g)
        y = ys + (i // 5) * (rh + g)
        card(s, x, y, cw, rh)
        txt(s, t, x + 0.08, y + 0.06, cw - 0.16, 0.32, sz=8, b=True, c=ACC[i % 6])
        txt(s, d, x + 0.08, y + 0.42, cw - 0.16, 0.40, sz=7, c=MUTED_FG)
    txt(
        s,
        "http://localhost:3001 · admin / harness · Anonymous viewer · Prometheus auto-provisioned",
        0.45,
        5.05,
        9.10,
        0.20,
        sz=8,
        c=MUTED_FG,
    )


def s8(p):
    s = p.slides.add_slide(p.slide_layouts[6])
    chrome(
        s,
        8,
        TS,
        "Weekly Team Velocity Scorecard",
        "Run these PromQL queries weekly · Targets based on production harness usage",
    )
    tbl(
        s,
        [
            ["Metric", "PromQL", "Target"],
            ["Total agent runs", "sum(harness_agent_runs_total)", "Rising w/w"],
            [
                "Runs per member",
                "sum by (user) (harness_agent_runs_total)",
                "Even dist.",
            ],
            [
                "Success rate",
                'sum(...{exit="ok"}) / sum(harness_agent_runs_total)',
                "> 90%",
            ],
            ["Generator success", '...{agent="generator",exit="ok"} / ...', "> 85%"],
            ["Evaluator pass", '...{agent="evaluator",exit="ok"} / ...', "> 80%"],
            [
                "Critic/generator",
                '...{agent="design-critic"} / ...{agent="generator"}',
                "< 2.0",
            ],
            ["Pending reviews", "harness_pending_reviews", "< 5/user"],
            [
                "Lane distribution",
                "sum by (lane) (harness_agent_runs_total)",
                "Bulk auto",
            ],
            ["Avg iterations", "avg(harness_iteration_current)", "< 5"],
            [
                "USD cost/day",
                "sum(max_over_time(claude_code_cost_usage_USD_total[24h]))",
                "Down/feat",
            ],
        ],
        0.45,
        1.20,
        9.10,
        3.70,
        [2.5, 4.5, 2.0],
    )
    txt(
        s,
        "(...) = harness_agent_runs_total. Paste into :9090 or Grafana.",
        0.45,
        5.00,
        9.10,
        0.25,
        sz=9,
        c=MUTED_FG,
    )


def s9(p):
    s = p.slides.add_slide(p.slide_layouts[6])
    chrome(
        s,
        9,
        TS,
        "Greenfield Metrics",
        "/scaffold -> /brd -> /spec -> /design -> /auto · All queryable from Prometheus :9090",
    )
    tbl(
        s,
        [
            ["Metric", "Source", "PromQL / Measurement"],
            [
                "Time-to-first-green",
                "Harness",
                "ts delta: first record -> evaluator exit=ok",
            ],
            [
                "Iterations-to-green",
                "Harness",
                "max(harness_iteration_current) per group",
            ],
            [
                "Pass-on-first-try",
                "Harness",
                '...{agent="evaluator",iteration="1",exit="ok"} / groups',
            ],
            [
                "Tokens per feature",
                "Native",
                "sum(increase(token_usage_total[24h])) / features",
            ],
            [
                "Cost per feature",
                "Native",
                "sum(increase(cost_usage_total[24h])) / features",
            ],
            [
                "Design-critic iters",
                "Harness",
                'sum(...{agent="design-critic"}) per group',
            ],
            [
                "Scaffold-to-first-PR",
                "External",
                "First harness ts -> PR merge ts via GitHub API",
            ],
        ],
        0.45,
        1.20,
        9.10,
        3.50,
        [2.5, 1.2, 6.3],
    )
    txt(
        s,
        'All harness metrics in Prometheus. Per-user: add {user="name"} filter.',
        0.45,
        4.80,
        9.10,
        0.35,
        sz=9,
        c=MUTED_FG,
    )


def s10(p):
    s = p.slides.add_slide(p.slide_layouts[6])
    chrome(
        s,
        10,
        TS,
        "Brownfield + Review ROI",
        "Attempt-to-merge > throughput · Review displacement = largest cost savings",
    )
    card(s, 0.45, 1.20, 4.45, 2.00)
    txt(s, "BROWNFIELD", 0.60, 1.28, 4.10, 0.25, sz=12, b=True, c=ACC[1])
    txt(
        s,
        "Attempt-to-merge — merged / opened PRs\nRework rate (7/30-day) — follow-up fixes\nBlast radius — LOC / files per PR\nSeam fit — changes inside /seam-finder output\nLane correctness — right lane for the work\nMap staleness — 14 days OR 50 commits",
        0.60,
        1.58,
        4.10,
        1.50,
        sz=9,
        c=BODY_FG,
    )
    card(s, 5.10, 1.20, 4.45, 2.00)
    txt(s, "REVIEW & SECURITY ROI", 5.25, 1.28, 4.10, 0.25, sz=12, b=True, c=ACC[2])
    txt(
        s,
        "Reviewer-hours per PR (pre vs post harness)\nPre-merge defects caught by agent reviewers\nReview agent cost — claude_code.cost.usage\nSecurity findings shifted left (pre-merge)\nReviewer-comments-per-PR (quality proxy)\nVerification tax — human / agent wall-clock",
        5.25,
        1.58,
        4.10,
        1.50,
        sz=9,
        c=BODY_FG,
    )
    call(
        s,
        0.45,
        3.45,
        9.10,
        0.70,
        "Monthly savings = PRs/mo x (baseline - post-harness review hrs) x loaded rate - agent cost\nAnchor: 200 PRs x 1.0 hr saved x $150/hr = $30k/month",
    )
    card(s, 0.45, 4.30, 4.45, 0.50, f=OK_BG)
    txt(
        s,
        "Yield > Throughput. Attempt-to-merge, rework rate, verification tax survive Goodhart.",
        0.60,
        4.38,
        4.10,
        0.35,
        sz=9,
        b=True,
        c=OK_FG,
    )
    card(s, 5.10, 4.30, 4.45, 0.50, f=ERR_BG)
    txt(
        s,
        "Quality drift: lint/complexity rises ~18-39% in AI-heavy repos. Watch defect-escape/1k LOC.",
        5.25,
        4.38,
        4.10,
        0.35,
        sz=9,
        b=True,
        c=ERR_FG,
    )


def s11(p):
    s = p.slides.add_slide(p.slide_layouts[6])
    chrome(
        s,
        11,
        TS,
        "Getting Started: 5 Minutes to Team Telemetry",
        "One-time per team · Per-developer: just /scaffold · Everything auto-configured",
    )
    steps = [
        (
            "1",
            "Clone & Load Plugin",
            "git clone .../claude_harness_eng_v4.git\nclaude --plugin-dir ~/claude_harness_eng_v4/.claude",
            ACC[0],
        ),
        (
            "2",
            "Run /scaffold",
            "Generates settings.json with telemetry env vars, .env, docker compose, Grafana configs",
            ACC[1],
        ),
        (
            "3",
            "Start Telemetry Stack",
            "docker compose -f telemetry_docker_compose.yml up -d\n4 services: Collector · Prometheus · Pushgateway · Grafana",
            ACC[2],
        ),
        (
            "4",
            "Open Grafana",
            "http://localhost:3001 (admin / harness)\n25-panel Team Productivity dashboard auto-loaded",
            ACC[3],
        ),
        (
            "5",
            "Build Something",
            "/brd -> /spec -> /design -> /auto\nMetrics flow automatically. Check Grafana.",
            ACC[4],
        ),
    ]
    for i, (n, t, d, a) in enumerate(steps):
        x, y = 0.45, 1.25 + i * 0.78
        card(s, x, y, 0.50, 0.65)
        txt(
            s,
            n,
            x + 0.10,
            y + 0.10,
            0.30,
            0.45,
            sz=22,
            b=True,
            c=a,
            a=PP_ALIGN.CENTER,
            v=MSO_ANCHOR.MIDDLE,
        )
        txt(s, t, 1.10, y + 0.05, 2.50, 0.30, sz=13, b=True, c=a)
        txt(s, d, 1.10, y + 0.32, 8.00, 0.35, sz=9, c=BODY_FG)
    txt(
        s,
        "Remote: change OTEL_EXPORTER_OTLP_ENDPOINT + HARNESS_PUSHGATEWAY_URL in settings.json.",
        0.45,
        5.00,
        9.10,
        0.25,
        sz=9,
        c=MUTED_FG,
    )


ALL = [s7, s8, s9, s10, s11]
