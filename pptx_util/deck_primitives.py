"""Shared primitives for PPTX deck generation — colors, helpers, layout constants."""
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BG = RGBColor(0x0F, 0x17, 0x2A)
ACCENT_BAR = RGBColor(0x38, 0xBD, 0xF8)
CARD = RGBColor(0x1E, 0x29, 0x3B)
TITLE_FG = RGBColor(0xF8, 0xFA, 0xFC)
SUBTITLE_FG = RGBColor(0x38, 0xBD, 0xF8)
BODY_FG = RGBColor(0xCB, 0xD5, 0xE1)
MUTED_FG = RGBColor(0x94, 0xA3, 0xB8)
TH_BG = RGBColor(0x33, 0x41, 0x55)
TA_BG = RGBColor(0x16, 0x20, 0x33)
TR_BG = RGBColor(0x1E, 0x29, 0x3B)
CALL_BG = RGBColor(0x0C, 0x4A, 0x6E)
OK_BG = RGBColor(0x1A, 0x36, 0x1A)
ERR_BG = RGBColor(0x36, 0x1A, 0x1A)
OK_FG = RGBColor(0x4A, 0xDE, 0x80)
ERR_FG = RGBColor(0xF8, 0x71, 0x71)
ACC = [
    RGBColor(0x38, 0xBD, 0xF8), RGBColor(0xFB, 0x92, 0x3C),
    RGBColor(0x4A, 0xDE, 0x80), RGBColor(0xA7, 0x8B, 0xFA),
    RGBColor(0xFB, 0xBF, 0x24), RGBColor(0xF8, 0x71, 0x71),
]
F = "Calibri"
SW, SH = 10.0, 5.625
BRAND = "Claude Harness Engine v4   ·   Telemetry & Team Velocity"

GR = [(0.45, 1.30), (3.50, 1.30), (6.55, 1.30), (0.45, 3.05), (3.50, 3.05), (6.55, 3.05)]
CW, CH = 2.95, 1.65

_TOTAL_SLIDES = None


def set_total_slides(n: int) -> None:
    global _TOTAL_SLIDES
    _TOTAL_SLIDES = n


def bg(s):
    b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(SW), Inches(SH))
    b.fill.solid(); b.fill.fore_color.rgb = BG; b.line.fill.background()
    t = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(SW), Inches(0.06))
    t.fill.solid(); t.fill.fore_color.rgb = ACCENT_BAR; t.line.fill.background()


def txt(s, t, x, y, w, h, sz=12, b=False, c=BODY_FG, a=PP_ALIGN.LEFT, v=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = v
    for i, ln in enumerate(t.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph(); p.alignment = a
        r = p.add_run(); r.text = ln; r.font.name = F; r.font.size = Pt(sz)
        r.font.bold = b; r.font.color.rgb = c
    return tb


def card(s, x, y, w, h, f=CARD):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = f; sh.line.fill.background()
    sh.adjustments[0] = 0.06; return sh


def chrome(s, n, ts=None, title="", sub=""):
    bg(s); txt(s, title, 0.45, 0.20, 9.10, 0.55, sz=26, b=True, c=TITLE_FG)
    txt(s, sub, 0.45, 0.78, 9.10, 0.35, sz=12, c=SUBTITLE_FG)
    txt(s, BRAND, 0.45, 5.30, 6.50, 0.20, sz=8, c=MUTED_FG)
    total = _TOTAL_SLIDES if _TOTAL_SLIDES is not None else ts
    if total is None:
        raise RuntimeError("call set_total_slides() or pass ts= to chrome()")
    txt(s, f"{n} / {total}", 8.50, 5.30, 1.10, 0.20, sz=8, c=MUTED_FG, a=PP_ALIGN.RIGHT)


def cardt(s, x, y, w, h, hd, bd, ac):
    card(s, x, y, w, h); txt(s, hd, x + .18, y + .12, w - .36, .30, sz=13, b=True, c=ac)
    txt(s, bd, x + .18, y + .45, w - .36, h - .55, sz=10, c=BODY_FG)


def six(s, cards):
    for i, (h, b) in enumerate(cards):
        cardt(s, *GR[i], CW, CH, h, b, ACC[i % 6])


def tbl(s, rows, x, y, w, h, cw):
    nr, nc = len(rows), len(rows[0])
    ts = s.shapes.add_table(nr, nc, Inches(x), Inches(y), Inches(w), Inches(h)).table
    tot = sum(cw)
    for ci, ratio in enumerate(cw): ts.columns[ci].width = Inches(w * ratio / tot)
    for ri, row in enumerate(rows):
        for ci, v in enumerate(row):
            c = ts.cell(ri, ci); c.margin_left = Inches(.08); c.margin_right = Inches(.08)
            c.margin_top = Inches(.04); c.margin_bottom = Inches(.04)
            c.vertical_anchor = MSO_ANCHOR.TOP; ih = ri == 0; c.fill.solid()
            c.fill.fore_color.rgb = TH_BG if ih else (TA_BG if ri % 2 == 0 else TR_BG)
            tf = c.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT; r = p.add_run(); r.text = v; r.font.name = F
            r.font.size = Pt(11 if ih else 9); r.font.bold = ih
            r.font.color.rgb = TITLE_FG if ih else BODY_FG


def call(s, x, y, w, h, t, cl=ACCENT_BAR):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = CALL_BG
    sh.line.color.rgb = cl; sh.line.width = Pt(1); sh.adjustments[0] = 0.10
    txt(s, t, x + .20, y + .10, w - .40, h - .20, sz=11, b=True, c=TITLE_FG, v=MSO_ANCHOR.MIDDLE)


def fbox(s, x, y, w, h, l, f):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = f; sh.line.fill.background()
    sh.adjustments[0] = 0.12
    txt(s, l, x + .08, y + .04, w - .16, h - .08, sz=9, b=True, c=TITLE_FG, a=PP_ALIGN.CENTER, v=MSO_ANCHOR.MIDDLE)


def arr(s, x, y, w, h, c=ACCENT_BAR):
    sh = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background()


def darr(s, x, y, w, h, c=ACCENT_BAR):
    sh = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = c; sh.line.fill.background()
