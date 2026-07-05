#!/usr/bin/env python3
"""Test helper for pe-ic-memo (not a public entry point): builds a small
sample deck via build_deck, re-opens the real generated file with
python-pptx, and prints a JSON summary for test/pe-ic-memo-render.test.js
to assert against.

Usage: python3 _test_build_sample.py <out_path>
"""

import json
import os
import sys

from pptx import Presentation

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from render_deck import build_deck  # noqa: E402


def _table_dims(slide):
    for shape in slide.shapes:
        if getattr(shape, "has_table", False):
            table = shape.table
            return {"rows": len(table.rows), "cols": len(table.columns)}
    return None


def _build_sample_memo():
    return {
        "title": {"company": "Acme Corp", "deal_name": "Project Falcon", "date": "2026-07-05"},
        "sections": [
            {
                "heading": "Executive Summary",
                "bullets": ["Deal rationale", "Recommendation: Proceed"],
                "table": None,
            },
            {
                "heading": "Financial Analysis",
                "bullets": ["Revenue grew 22% YoY"],
                "table": {
                    "headers": ["Year", "Revenue", "EBITDA"],
                    "rows": [["2024", "$50M", "$10M"], ["2025", "$61M", "$14M"]],
                },
            },
        ],
    }


def _summarize_deck(out_path):
    prs = Presentation(out_path)
    slides = list(prs.slides)
    return {
        "slide_count": len(slides),
        "titles": [s.shapes.title.text if s.shapes.title else None for s in slides],
        "table_dims_by_slide": [_table_dims(s) for s in slides],
    }


def main():
    out_path = sys.argv[1]
    memo = _build_sample_memo()
    build_deck(memo, out_path)
    summary = _summarize_deck(out_path)
    print(json.dumps(summary))


if __name__ == "__main__":
    main()
