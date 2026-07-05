"""Render a PE IC memo as a PPTX deck.

memo shape:
{
    "title": {"company": str, "deal_name": str, "date": str},
    "sections": [
        {
            "heading": str,
            "bullets": [str, ...],
            "table": {"headers": [str, ...], "rows": [[str, ...], ...]} | None,
        },
        ...
    ],
}

Renders a generic layout using python-pptx's default template: a title slide
followed by one content slide per section. Sections with a `table` also get
a native PPTX table shape. See this skill's SKILL.md "Firm Branding" section
for how a real firm template gets swapped in later -- that swap is not
implemented here.
"""

import os

from pptx import Presentation
from pptx.util import Inches


def _add_title_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title["company"]
    subtitle = slide.placeholders[1]
    subtitle.text = "{} — CONFIDENTIAL — Investment Committee\n{}".format(
        title["deal_name"], title["date"]
    )


def _add_bullets(body_placeholder, bullets):
    tf = body_placeholder.text_frame
    tf.clear()
    tf.text = bullets[0]
    for bullet in bullets[1:]:
        p = tf.add_paragraph()
        p.text = bullet


def _add_table(slide, table_data, top):
    rows = len(table_data["rows"]) + 1
    cols = len(table_data["headers"])
    left = Inches(0.5)
    width = Inches(12.33)
    height = Inches(0.4 * rows)
    shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = shape.table
    for c, header in enumerate(table_data["headers"]):
        table.cell(0, c).text = header
    for r, row in enumerate(table_data["rows"]):
        for c, value in enumerate(row):
            table.cell(r + 1, c).text = str(value)


def _add_content_slide(prs, section):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = section["heading"]
    body = slide.placeholders[1]
    if section.get("bullets"):
        _add_bullets(body, section["bullets"])
    else:
        body.text_frame.clear()
    if section.get("table"):
        top = Inches(4.2) if section.get("bullets") else Inches(1.8)
        _add_table(slide, section["table"], top)
    return slide


def build_deck(memo, out_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _add_title_slide(prs, memo["title"])
    for section in memo["sections"]:
        _add_content_slide(prs, section)

    out_dir = os.path.dirname(out_path)
    if out_dir:
        os.makedirs(out_dir, exist_ok=True)
    prs.save(out_path)
    return out_path
